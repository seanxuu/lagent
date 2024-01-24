from typing import Dict, Optional, Type

from pptx import Presentation

from lagent.actions.base_action import BaseAction
from lagent.actions.parser import BaseParser, JsonParser

DEFAULT_DESCRIPTION = dict(
    name='PPT',
    description=
    'This tool allows you to create ppt slides with text, paragraph, images, with good looking styles',
    api_list=[
        dict(
            name='create_file',
            description='Create a pptx file with specific themes',
            parameters=[
                dict(
                    name='theme', type='STRING', description='the theme used'),
                dict(
                    name='abs_location',
                    type='STRING',
                    description='the ppt file\'s absolute location')
            ],
            required=['theme', 'abs_location'],
            return_data=[
                dict(name='status', description='the result of the execution')
            ]),
        dict(
            name='get_image',
            description=
            'Get an image given comma separated keywords, return the image path.',
            parameters=[
                dict(
                    name='keywords',
                    type='STRING',
                    description=
                    'the comma separated keywords to describe the image')
            ],
            required=['keywords'],
            return_data=[
                dict(name='status', description='the result of the execution')
            ]),
        dict(
            name='add_first_page',
            description='Add the first page of ppt.',
            parameters=[
                dict(
                    name='title',
                    type='STRING',
                    description='the title of ppt'),
                dict(
                    name='subtitle',
                    type='STRING',
                    description='the subtitle of ppt')
            ],
            required=['title', 'subtitle'],
            return_data=[
                dict(name='status', description='the result of the execution')
            ]),
        dict(
            name='add_text_page',
            description='Add text page of ppt',
            parameters=[
                dict(
                    name='title',
                    type='STRING',
                    description='the title of the page'),
                dict(
                    name='bullet_items',
                    type='STRING',
                    description=
                    'bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.'
                )
            ],
            required=['title', 'bullet_items'],
            return_data=[
                dict(name='status', description='the result of the execution')
            ]),
        dict(
            name='add_text_image_page',
            description=
            'Add a text page with one image. Image should be a path',
            parameters=[
                dict(
                    name='title',
                    type='STRING',
                    description='the title of the page'),
                dict(
                    name='bullet_items',
                    type='STRING',
                    description=
                    'bullet_items should be string, for multiple bullet items, please use [SPAN] to separate them.'
                ),
                dict(
                    name='image',
                    type='STRING',
                    description='the path of the image')
            ],
            required=['title', 'bullet_items', 'image'],
            return_data=[
                dict(name='status', description='the result of the execution')
            ]),
        dict(
            name='submit_file',
            description=
            'When all steps done, YOU MUST use submit_file() to submit your work.',
            parameters=[],
            required=[],
            return_data=[
                dict(name='status', description='the result of the execution')
            ])
    ])

THEME_MAPPING = {
    'Default': {
        'template': None,
        'title': 'Title Slide',
        'single': 'Title and Content',
        'two': 'Tow content',
    }
}


class PPT(BaseAction):
    """Plugin to create ppt slides with text, paragraph, images in good looking styles"""

    def __init__(self,
                 theme_mapping: Optional[Dict[str, dict]] = None,
                 description: Optional[dict] = None,
                 parser: Type[BaseParser] = JsonParser,
                 enable: bool = True):
        super().__init__(description or DEFAULT_DESCRIPTION, parser, enable)
        self.theme_mapping = theme_mapping or THEME_MAPPING
        self.pointer = None
        self.location = None

    def create_file(self, theme: str, abs_location: str) -> dict:
        self.location = abs_location
        try:
            self.pointer = Presentation(self.theme_mapping[theme]['template'])
            self.pointer.slide_master.name = theme
            # print('created')
        except Exception as e:
            print(e)
        return dict(status='created a ppt file.')

    def add_first_page(self, title: str, subtitle: str) -> dict:
        layout_name = self.theme_mapping[
            self.pointer.slide_master.name]['title']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_subtitle = slide.placeholders
        ph_title.text = title
        if subtitle:
            ph_subtitle.text = subtitle
        return dict(status='added page')

    def add_text_page(self, title: str, bullet_items: str) -> dict:
        layout_name = self.theme_mapping[
            self.pointer.slide_master.name]['single']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_body = slide.placeholders
        ph_title.text = title
        ph = ph_body
        tf = ph.text_frame
        for i, item in enumerate(bullet_items.split('[SPAN]')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item.strip()
            p.level = 0
        return dict(status='added page')

    def add_text_image_page(self, title: str, bullet_items: str,
                            image: str) -> dict:
        layout_name = self.theme_mapping[self.pointer.slide_master.name]['two']
        layout = next(i for i in self.pointer.slide_master.slide_layouts
                      if i.name == layout_name)
        slide = self.pointer.slides.add_slide(layout)
        ph_title, ph_body1, ph_body2 = slide.placeholders
        ph_title.text = title
        ph = ph_body2
        image_pil = image.to_pil()
        left = ph.left
        width = ph.width
        height = int(width / image_pil.width * image_pil.height)
        top = (ph.top + (ph.top + ph.height)) // 2 - height // 2
        slide.shapes.add_picture(image.to_path(), left, top, width, height)

        ph = ph_body1
        tf = ph.text_frame
        for i, item in enumerate(bullet_items.split('[SPAN]')):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item.strip()
            p.level = 0

        return dict(status='added page')

    def submit_file(self) -> dict:
        # file_path = os.path.join(self.CACHE_DIR, f'{self._return_timestamp()}.pptx')
        # self.pointer.save(file_path)
        # retreival_url = upload_file(file_path)
        self.pointer.save(self.location)
        return dict(status=f'submitted. view ppt at {self.location}')